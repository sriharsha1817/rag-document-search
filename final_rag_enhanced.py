import streamlit as st
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from transformers import pipeline, AutoTokenizer, AutoModelForCausalLM
from pydantic import BaseModel, Field
from typing import List
from pypdf import PdfReader
import pptx
import torch

st.set_page_config(page_title="AI Document Search", layout="wide")

class RAGAnswer(BaseModel):
    answer: str = Field(description="Answer to the question")
    confidence: str = Field(description="Confidence: high/medium/low")
    source_chunks: List[str] = Field(description="Top 3 relevant text chunks")

@st.cache_resource
def load_models():
    # Load embedding model
    embeddings = SentenceTransformer('all-MiniLM-L6-v2')
    
    # Load a small but capable local LLM
    try:
        # Use a lightweight model that works well for Q&A
        qa_pipeline = pipeline(
            "text-generation",
            model="microsoft/DialoGPT-small",
            tokenizer="microsoft/DialoGPT-small",
            device=0 if torch.cuda.is_available() else -1
        )
        return embeddings, qa_pipeline
    except:
        # Fallback to even simpler model
        qa_pipeline = pipeline(
            "question-answering",
            model="distilbert-base-cased-distilled-squad"
        )
        return embeddings, qa_pipeline

embeddings, qa_model = load_models()

def process_documents(text):
    chunks = [para.strip() for para in text.split("\n\n") if para.strip()]
    return chunks

def generate_answer_with_local_llm(context, query, qa_model):
    """Generate answer using local LLM"""
    try:
        # If it's a question-answering pipeline
        if hasattr(qa_model, 'task') and qa_model.task == 'question-answering':
            result = qa_model(question=query, context=context)
            answer = result['answer']
            confidence = "high" if result['score'] > 0.7 else "medium" if result['score'] > 0.3 else "low"
        else:
            # If it's a text generation pipeline
            prompt = f"Context: {context[:1000]}\n\nQuestion: {query}\n\nAnswer:"
            result = qa_model(prompt, max_length=200, num_return_sequences=1, temperature=0.7)
            answer = result[0]['generated_text'].split("Answer:")[-1].strip()
            confidence = "medium"
        
        return answer, confidence
    except Exception as e:
        st.error(f"LLM Error: {e}")
        return "I couldn't generate an answer due to a processing error.", "low"

def rag_query(chunks, query, top_k=3):
    try:
        # Use local embeddings for similarity search
        doc_embeddings = embeddings.encode(chunks)
        query_embedding = embeddings.encode([query])
        scores = cosine_similarity(query_embedding, doc_embeddings)[0]
        
        # Get top relevant chunks
        top_indices = scores.argsort()[-top_k:][::-1]
        source_chunks = [chunks[i] for i in top_indices]
        context = "\n\n".join(source_chunks)
        
        # Generate answer using local LLM
        answer, confidence = generate_answer_with_local_llm(context, query, qa_model)
        
        response = RAGAnswer(
            answer=answer,
            confidence=confidence,
            source_chunks=source_chunks
        )
        
        return context, response
    
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return "", RAGAnswer(
            answer="An error occurred while processing your query.", 
            confidence="low", 
            source_chunks=[]
        )

st.title("🤖 Enhanced Local AI Document Search")
st.markdown("Upload documents → Chat with local AI models (No API keys needed!)")

# Show model info
st.sidebar.info("🧠 Using local models:\n- Embeddings: all-MiniLM-L6-v2\n- LLM: DistilBERT/DialoGPT")

uploaded_files = st.sidebar.file_uploader(
    "Choose TXT/PDF/PPT files", accept_multiple_files=True, 
    type=['txt', 'pdf', 'pptx']
)

if uploaded_files:
    documents_text = ""
    for file in uploaded_files:
        if file.type == "text/plain":
            documents_text += file.read().decode("utf-8") + "\n\n"
        elif file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                documents_text += page.extract_text() + "\n\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents_text += shape.text + "\n\n"
    
    st.sidebar.success(f"✅ Loaded {len(process_documents(documents_text))} chunks")
    st.session_state.documents = documents_text

if "messages" not in st.session_state:
    st.session_state.messages = []

if "documents" in st.session_state:
    st.header("💬 Chat with Documents")
    
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    if prompt := st.chat_input("Ask about your documents..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("🔍 Processing with local AI..."):
                context, answer = rag_query(
                    process_documents(st.session_state.documents), 
                    prompt
                )
                
                st.markdown(f"**Answer:** {answer.answer}")
                st.markdown(f"**Confidence:** {answer.confidence}")
                
                with st.expander("📖 View Sources", expanded=False):
                    for i, chunk in enumerate(answer.source_chunks, 1):
                        st.markdown(f"**Source {i}:**")
                        st.markdown(chunk)

        st.session_state.messages.append({"role": "assistant", "content": answer.answer})

st.markdown("---")
st.markdown("**✅ Features:** RAG • Local Embeddings • Local LLM • 100% Free • Better Accuracy")