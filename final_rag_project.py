import streamlit as st
from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI, GoogleGenerativeAIEmbeddings
from sklearn.metrics.pairwise import cosine_similarity
import os
from langchain_core.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field
from typing import List
from pypdf import PdfReader
import pptx
import time


# First Streamlit command
st.set_page_config(page_title="AI Document Search", layout="wide")

load_dotenv()

# Pydantic model for structured output (from ma'am's PPT)
class RAGAnswer(BaseModel):
    answer: str = Field(description="Answer to the question")
    confidence: str = Field(description="Confidence: high/medium/low")
    source_chunks: List[str] = Field(description="Top 3 relevant text chunks")

@st.cache_resource
def load_models():
    llm = ChatGoogleGenerativeAI(model="gemini-1.5-flash")
    embeddings = GoogleGenerativeAIEmbeddings(model="models/text-embedding-004")
    parser = PydanticOutputParser(pydantic_object=RAGAnswer)
    return llm, embeddings, parser

llm, embeddings, parser = load_models()

def process_documents(text):
    """Split text into chunks for RAG"""
    chunks = [para.strip() for para in text.split("\n\n") if para.strip()]
    return chunks

def rag_query(chunks, query, top_k=3):
    """RAG with Pydantic structured output"""
    try:
        doc_embeddings = embeddings.embed_documents(chunks)
        query_embedding = embeddings.embed_query(query)
        scores = cosine_similarity([query_embedding], doc_embeddings)[0]
        
        # Get top K chunks with indices
        top_indices = scores.argsort()[-top_k:][::-1]
        source_chunks = [chunks[i] for i in top_indices]
        context = "\n\n".join(source_chunks)
        
        prompt = f"""Answer based ONLY on the CONTEXT. Return JSON with confidence and sources.

CONTEXT:
{context}

QUESTION: {query}

SCHEMA:
{parser.get_format_instructions()}"""
        
        # Add delay to avoid rate limiting
        time.sleep(2)
        
        # Use Runnable chain: llm → parser
        chain = llm | parser
        response = chain.invoke(prompt)
        return context, response
    
    except Exception as e:
        if "429" in str(e) or "Quota exceeded" in str(e):
            st.error("⚠️ Rate limit exceeded. Please wait a moment and try again.")
            return "", RAGAnswer(answer="Rate limit exceeded. Please try again in a few minutes.", confidence="low", source_chunks=[])
        else:
            st.error(f"Error: {str(e)}")
            return "", RAGAnswer(answer="An error occurred while processing your query.", confidence="low", source_chunks=[])

# UI
st.title("🤖 AI Document Search & Knowledge Retrieval")
st.markdown("Upload documents → Chat with structured RAG answers + sources!")



# Install: pip install pypdf python-pptx
from pypdf import PdfReader
import pptx

uploaded_files = st.sidebar.file_uploader(
    "Choose TXT/PDF/PPT files", accept_multiple_files=True, 
    type=['txt', 'pdf', 'pptx']
)

if uploaded_files:
    documents_text = ""
    for file in uploaded_files:
        if file.type == "text/plain":
            documents_text += file.read().decode("utf-8") + "\n\n"
        elif file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                documents_text += page.extract_text() + "\n\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents_text += shape.text + "\n\n"
    
    st.sidebar.success(f"✅ Loaded {len(process_documents(documents_text))} chunks")
    st.session_state.documents = documents_text


# Chat interface
if "messages" not in st.session_state:
    st.session_state.messages = []

if "documents" in st.session_state:
    st.header("💬 Chat with Documents")
    
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    if prompt := st.chat_input("Ask about your documents..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("🔍 RAG Search..."):
                context, answer = rag_query(
                    process_documents(st.session_state.documents), 
                    prompt
                )
                
                # Display structured answer
                st.markdown(f"**Answer:** {answer.answer}")
                st.markdown(f"**Confidence:** {answer.confidence}")
                
                with st.expander("📖 View Sources", expanded=False):
                    for i, chunk in enumerate(answer.source_chunks, 1):
                        st.markdown(f"**Source {i}:**")
                        st.markdown(chunk)

        st.session_state.messages.append({"role": "assistant", "content": str(answer)})

st.markdown("---")
st.markdown("**✅ Features:** RAG • Pydantic Parser • Source Attribution • Confidence Scoring")
