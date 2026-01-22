import streamlit as st
from dotenv import load_dotenv
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from sklearn.metrics.pairwise import cosine_similarity
import os
from langchain_core.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field
from typing import List
from pypdf import PdfReader
import pptx

st.set_page_config(page_title="AI Document Search", layout="wide")

load_dotenv()

class RAGAnswer(BaseModel):
    answer: str = Field(description="Answer to the question")
    confidence: str = Field(description="Confidence: high/medium/low")
    source_chunks: List[str] = Field(description="Top 3 relevant text chunks")

@st.cache_resource
def load_models():
    # Try Streamlit secrets first, then environment variables
    api_key = st.secrets.get("OPENAI_API_KEY") or os.getenv("OPENAI_API_KEY")
    if not api_key:
        st.error("OpenAI API key not found. Please add it to Streamlit secrets or .env file.")
        return None, None, None
    
    llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0, api_key=api_key)
    embeddings = OpenAIEmbeddings(model="text-embedding-3-small", api_key=api_key)
    parser = PydanticOutputParser(pydantic_object=RAGAnswer)
    return llm, embeddings, parser

llm, embeddings, parser = load_models()

# Check if models loaded successfully
if llm is None:
    st.stop()

def process_documents(text):
    chunks = [para.strip() for para in text.split("\n\n") if para.strip()]
    return chunks

def rag_query(chunks, query, top_k=3):
    try:
        doc_embeddings = embeddings.embed_documents(chunks)
        query_embedding = embeddings.embed_query(query)
        scores = cosine_similarity([query_embedding], doc_embeddings)[0]
        
        top_indices = scores.argsort()[-top_k:][::-1]
        source_chunks = [chunks[i] for i in top_indices]
        context = "\n\n".join(source_chunks)
        
        prompt = f"""Answer based ONLY on the CONTEXT. Return JSON with confidence and sources.

CONTEXT:
{context}

QUESTION: {query}

SCHEMA:
{parser.get_format_instructions()}"""
        
        chain = llm | parser
        response = chain.invoke(prompt)
        return context, response
    
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return "", RAGAnswer(answer="An error occurred while processing your query.", confidence="low", source_chunks=[])

st.title("🤖 AI Document Search & Knowledge Retrieval")
st.markdown("Upload documents → Chat with structured RAG answers + sources!")

uploaded_files = st.sidebar.file_uploader(
    "Choose TXT/PDF/PPT files", accept_multiple_files=True, 
    type=['txt', 'pdf', 'pptx']
)

if uploaded_files:
    documents_text = ""
    for file in uploaded_files:
        if file.type == "text/plain":
            documents_text += file.read().decode("utf-8") + "\n\n"
        elif file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                documents_text += page.extract_text() + "\n\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents_text += shape.text + "\n\n"
    
    st.sidebar.success(f"✅ Loaded {len(process_documents(documents_text))} chunks")
    st.session_state.documents = documents_text

if "messages" not in st.session_state:
    st.session_state.messages = []

if "documents" in st.session_state:
    st.header("💬 Chat with Documents")
    
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    if prompt := st.chat_input("Ask about your documents..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("🔍 RAG Search..."):
                context, answer = rag_query(
                    process_documents(st.session_state.documents), 
                    prompt
                )
                
                st.markdown(f"**Answer:** {answer.answer}")
                st.markdown(f"**Confidence:** {answer.confidence}")
                
                with st.expander("📖 View Sources", expanded=False):
                    for i, chunk in enumerate(answer.source_chunks, 1):
                        st.markdown(f"**Source {i}:**")
                        st.markdown(chunk)

        st.session_state.messages.append({"role": "assistant", "content": str(answer)})

st.markdown("---")
st.markdown("**✅ Features:** RAG • OpenAI • Source Attribution • Confidence Scoring")