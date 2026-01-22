import streamlit as st
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from pydantic import BaseModel, Field
from typing import List
from pypdf import PdfReader
import pptx
import json

st.set_page_config(page_title="AI Document Search", layout="wide")

class RAGAnswer(BaseModel):
    answer: str = Field(description="Answer to the question")
    confidence: str = Field(description="Confidence: high/medium/low")
    source_chunks: List[str] = Field(description="Top 3 relevant text chunks")

@st.cache_resource
def load_embedding_model():
    return SentenceTransformer('all-MiniLM-L6-v2')

embeddings = load_embedding_model()

def process_documents(text):
    chunks = [para.strip() for para in text.split("\n\n") if para.strip()]
    return chunks

def simple_rag_answer(context, query):
    """Simple rule-based answer generation"""
    context_lower = context.lower()
    query_lower = query.lower()
    
    # Simple keyword matching and extraction
    sentences = context.split('.')
    relevant_sentences = []
    
    query_words = query_lower.split()
    for sentence in sentences:
        sentence_lower = sentence.lower()
        matches = sum(1 for word in query_words if word in sentence_lower)
        if matches > 0:
            relevant_sentences.append((sentence.strip(), matches))
    
    # Sort by relevance
    relevant_sentences.sort(key=lambda x: x[1], reverse=True)
    
    if relevant_sentences:
        answer = ". ".join([sent[0] for sent in relevant_sentences[:3]])
        confidence = "high" if relevant_sentences[0][1] >= 2 else "medium"
    else:
        answer = "I couldn't find specific information about your question in the provided documents."
        confidence = "low"
    
    return answer, confidence

def rag_query(chunks, query, top_k=3):
    try:
        # Use local embeddings
        doc_embeddings = embeddings.encode(chunks)
        query_embedding = embeddings.encode([query])
        scores = cosine_similarity(query_embedding, doc_embeddings)[0]
        
        top_indices = scores.argsort()[-top_k:][::-1]
        source_chunks = [chunks[i] for i in top_indices]
        context = "\n\n".join(source_chunks)
        
        # Generate answer using simple rules
        answer, confidence = simple_rag_answer(context, query)
        
        response = RAGAnswer(
            answer=answer,
            confidence=confidence,
            source_chunks=source_chunks
        )
        
        return context, response
    
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return "", RAGAnswer(
            answer="An error occurred while processing your query.", 
            confidence="low", 
            source_chunks=[]
        )

st.title("🤖 Free AI Document Search & Knowledge Retrieval")
st.markdown("Upload documents → Chat with local AI (No API keys needed!)")

uploaded_files = st.sidebar.file_uploader(
    "Choose TXT/PDF/PPT files", accept_multiple_files=True, 
    type=['txt', 'pdf', 'pptx']
)

if uploaded_files:
    documents_text = ""
    for file in uploaded_files:
        if file.type == "text/plain":
            documents_text += file.read().decode("utf-8") + "\n\n"
        elif file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                documents_text += page.extract_text() + "\n\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents_text += shape.text + "\n\n"
    
    st.sidebar.success(f"✅ Loaded {len(process_documents(documents_text))} chunks")
    st.session_state.documents = documents_text

if "messages" not in st.session_state:
    st.session_state.messages = []

if "documents" in st.session_state:
    st.header("💬 Chat with Documents")
    
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    if prompt := st.chat_input("Ask about your documents..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("🔍 RAG Search..."):
                context, answer = rag_query(
                    process_documents(st.session_state.documents), 
                    prompt
                )
                
                st.markdown(f"**Answer:** {answer.answer}")
                st.markdown(f"**Confidence:** {answer.confidence}")
                
                with st.expander("📖 View Sources", expanded=False):
                    for i, chunk in enumerate(answer.source_chunks, 1):
                        st.markdown(f"**Source {i}:**")
                        st.markdown(chunk)

        st.session_state.messages.append({"role": "assistant", "content": answer.answer})

st.markdown("---")
st.markdown("**✅ Features:** RAG • Local Embeddings • Rule-based Answers • 100% Free")