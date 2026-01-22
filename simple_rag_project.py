import os
from dotenv import load_dotenv
import streamlit as st
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_core.output_parsers import PydanticOutputParser
from pydantic import BaseModel, Field
from typing import List
from pypdf import PdfReader
import pptx

load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GOOGLE_API_KEY:
    raise ValueError("GOOGLE_API_KEY not found. Please check your .env file.")

st.set_page_config(page_title="AI Document Search", layout="wide")

class RAGAnswer(BaseModel):
    answer: str = Field(description="Answer to the question")
    confidence: str = Field(description="Confidence: high/medium/low")
    source_chunks: List[str] = Field(description="Top 3 relevant text chunks")

@st.cache_resource
def load_models():
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash",
        google_api_key=GOOGLE_API_KEY
    )
    parser = PydanticOutputParser(pydantic_object=RAGAnswer)
    return llm, parser

llm, parser = load_models()

def process_documents(text):
    return [para.strip() for para in text.split("\n\n") if para.strip()]

def simple_rag_query(chunks, query, top_k=3):
    # Simple keyword matching instead of embeddings
    query_words = set(query.lower().split())
    
    chunk_scores = []
    for i, chunk in enumerate(chunks):
        chunk_words = set(chunk.lower().split())
        score = len(query_words.intersection(chunk_words))
        chunk_scores.append((score, i, chunk))
    
    # Get top chunks by keyword overlap
    chunk_scores.sort(reverse=True)
    top_chunks = [chunk for _, _, chunk in chunk_scores[:top_k]]
    
    context = "\n\n".join(top_chunks)
    
    prompt = f"""Answer based ONLY on the CONTEXT.
Return JSON strictly following the schema.

CONTEXT:
{context}

QUESTION:
{query}

SCHEMA:
{parser.get_format_instructions()}
"""
    
    chain = llm | parser
    response = chain.invoke(prompt)
    return context, response

st.title("🤖 AI Document Search & Knowledge Retrieval")
st.markdown("Upload documents → Chat with structured answers + sources!")

uploaded_files = st.sidebar.file_uploader(
    "Choose TXT/PDF/PPT files",
    accept_multiple_files=True,
    type=["txt", "pdf", "pptx"]
)

if uploaded_files:
    documents_text = ""
    
    for file in uploaded_files:
        if file.type == "text/plain":
            documents_text += file.read().decode("utf-8") + "\n\n"
        elif file.type == "application/pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    documents_text += text + "\n\n"
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            presentation = pptx.Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        documents_text += shape.text + "\n\n"
    
    chunks = process_documents(documents_text)
    st.sidebar.success(f"✅ Loaded {len(chunks)} chunks")
    st.session_state.documents = documents_text

if "messages" not in st.session_state:
    st.session_state.messages = []

if "documents" in st.session_state:
    st.header("💬 Chat with Documents")
    
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    if prompt := st.chat_input("Ask about your documents..."):
        st.session_state.messages.append({"role": "user", "content": prompt})
        
        with st.chat_message("user"):
            st.markdown(prompt)
        
        with st.chat_message("assistant"):
            with st.spinner("🔍 Searching..."):
                context, answer = simple_rag_query(
                    process_documents(st.session_state.documents),
                    prompt
                )
                
                st.markdown(f"**Answer:** {answer.answer}")
                st.markdown(f"**Confidence:** {answer.confidence}")
                
                with st.expander("📖 View Sources"):
                    for i, chunk in enumerate(answer.source_chunks, 1):
                        st.markdown(f"**Source {i}:**")
                        st.markdown(chunk)
        
        st.session_state.messages.append(
            {"role": "assistant", "content": answer.answer}
        )

st.markdown("---")
st.markdown("**✅ Features:** Simple RAG • Gemini • Keyword Matching • Source Attribution")